"""Format detection and single-file conversion to readable text.

Used by `gaz preview` (bounded, to a terminal) and `gaz convert`
(unbounded, to a file) — both call convert_to_text() and differ only in
how much of the result they keep and where it goes.

Converter priority per format, each step attempted only if the previous
one is unavailable:
    docx/pptx/xlsx -> pandoc -> matching optional Python library
    pdf             -> pdftotext -> pandoc (n/a, no pdf input) -> pypdf
    yaml            -> PyYAML (no stdlib option)
    toml            -> stdlib tomllib (3.11+) -> tomli
    json/xml/csv    -> stdlib (always available)
    md/txt/unknown  -> read as-is
"""

from __future__ import annotations

import csv
import importlib.util
import io
import json
import shutil
import subprocess
import sys
import xml.dom.minidom
import xml.parsers.expat
from dataclasses import dataclass

# gaz convert is binary -> text only (see DESIGN.md). These formats are
# already text, so pretty-printing them is preview's job; convert refuses
# them rather than silently no-op'ing or inventing a translation feature.
PREVIEW_ONLY_FORMATS = {"json", "yaml", "toml", "xml"}

# Output formats convert() will actually produce, keyed by requested --to
# (or inferred output extension). "text" covers .md/.txt/.markdown, which
# are treated as equivalent since most sources only ever produce prose.
_OUTPUT_FORMAT_ALIASES = {
    "md": "text",
    "markdown": "text",
    "txt": "text",
    "text": "text",
    "csv": "csv",
}

_EXT_TO_FORMAT = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".pdf": "pdf",
    ".json": "json",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".toml": "toml",
    ".xml": "xml",
    ".csv": "csv",
    ".md": "md",
    ".markdown": "md",
    ".txt": "txt",
}

# Formats pandoc's own format names for --from, when they differ from ours.
_PANDOC_FROM = {
    "docx": "docx",
    "pptx": "pptx",
    "xlsx": "xlsx",
}


@dataclass
class ConvertResult:
    text: str
    method: str
    complete: bool = True
    warning: str | None = None


class UnsupportedFormat(Exception):
    """No converter available for this input. Message names what's missing."""


def detect_format(path: str) -> str:
    """Extension-based format detection, falling back to a binary/text sniff."""
    for ext, fmt in _EXT_TO_FORMAT.items():
        if path.lower().endswith(ext):
            return fmt
    return _sniff_format(path)


def _sniff_format(path: str) -> str:
    try:
        with open(path, "rb") as f:
            chunk = f.read(8192)
    except OSError:
        return "unknown"
    if b"\x00" in chunk:
        return "binary"
    return "txt"


def _has_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def convert_to_text(
    path: str, *, max_seconds: float = 30.0, to_format: str | None = None
) -> ConvertResult:
    """Convert `path` to readable text, dispatching by detected format.

    `to_format` (one of "text"/"md"/"txt"/"csv", or None) is only meaningful
    for formats with more than one sensible text representation — currently
    just xlsx, which can render as a Markdown-ish table or as real CSV.
    Every other format ignores it and produces its one natural
    representation (pandoc/pdftotext output, a pretty-printer's output,
    etc.) regardless of what was requested; callers that care about the
    binary->text-only /  preview-only-formats-are-rejected distinction
    (i.e. `gaz convert`) should check `detect_format()` against
    PREVIEW_ONLY_FORMATS themselves before calling this, since preview_only
    formats are valid input for `gaz preview` but not for `gaz convert`.
    """
    fmt = detect_format(path)
    normalized_to = _OUTPUT_FORMAT_ALIASES.get((to_format or "").lower())

    if fmt == "xlsx" and normalized_to == "csv":
        return _convert_xlsx_to_csv(path, max_seconds=max_seconds)
    if fmt in ("docx", "pptx", "xlsx"):
        return _convert_office(path, fmt, max_seconds=max_seconds)
    if fmt in ("xls",):
        return _convert_xls(path, max_seconds=max_seconds)
    if fmt == "pdf":
        return _convert_pdf(path, max_seconds=max_seconds)
    if fmt == "json":
        return _pretty_json(path)
    if fmt == "yaml":
        return _pretty_yaml(path)
    if fmt == "toml":
        return _pretty_toml(path)
    if fmt == "xml":
        return _pretty_xml(path)
    if fmt == "csv":
        return _pretty_csv(path)
    if fmt in ("md", "txt"):
        return _read_as_is(path, method="text")
    if fmt == "binary":
        raise UnsupportedFormat(
            f"cannot preview {path!r} — file appears to be binary with no "
            f"recognized extension. No converter to try."
        )
    # Unknown extension, sniffed as text: just read it.
    return _read_as_is(path, method="text (unrecognized extension)")


def _read_as_is(path: str, *, method: str) -> ConvertResult:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return ConvertResult(text=f.read(), method=method)
    except OSError as e:
        raise UnsupportedFormat(f"cannot read {path!r}: {e}")


def _run_subprocess(cmd: list[str], *, max_seconds: float) -> subprocess.CompletedProcess:
    return subprocess.run(
        cmd,
        capture_output=True,
        timeout=max_seconds,
        check=False,
    )


def _convert_office(path: str, fmt: str, *, max_seconds: float) -> ConvertResult:
    if shutil.which("pandoc"):
        try:
            proc = _run_subprocess(
                ["pandoc", "--from", _PANDOC_FROM[fmt], "--to", "markdown", path],
                max_seconds=max_seconds,
            )
        except subprocess.TimeoutExpired:
            return ConvertResult(
                text="", method="pandoc", complete=False,
                warning=f"pandoc did not finish within {max_seconds}s",
            )
        if proc.returncode == 0:
            warning = proc.stderr.decode("utf-8", errors="replace").strip() or None
            return ConvertResult(
                text=proc.stdout.decode("utf-8", errors="replace"),
                method="pandoc",
                warning=warning,
            )
        # pandoc failed outright — fall through to a Python fallback below.

    if fmt == "docx" and _has_module("docx"):
        return _convert_docx_python(path)
    if fmt == "pptx" and _has_module("pptx"):
        return _convert_pptx_python(path)
    if fmt == "xlsx" and _has_module("openpyxl"):
        return _convert_xlsx_python(path)

    raise UnsupportedFormat(
        f"cannot convert .{fmt} — no converter available. "
        f"Install pandoc (https://pandoc.org/installing.html), or run "
        f"`pip install gaz[preview]` for a Python fallback."
    )


def _convert_docx_python(path: str) -> ConvertResult:
    import docx

    document = docx.Document(path)
    paragraphs = [p.text for p in document.paragraphs]
    return ConvertResult(text="\n\n".join(paragraphs), method="python-docx")


def _convert_pptx_python(path: str) -> ConvertResult:
    from pptx import Presentation

    presentation = Presentation(path)
    slides_text = []
    for i, slide in enumerate(presentation.slides, start=1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
        slides_text.append("\n".join(lines))
    return ConvertResult(text="\n\n".join(slides_text), method="python-pptx")


def _convert_xlsx_python(path: str) -> ConvertResult:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets_text = []
    for name in workbook.sheetnames:
        sheet = workbook[name]
        buf = io.StringIO()
        writer = csv.writer(buf)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow(["" if v is None else v for v in row])
        sheets_text.append(f"## Sheet: {name}\n\n{buf.getvalue()}")
    return ConvertResult(text="\n\n".join(sheets_text), method="openpyxl")


def _convert_xlsx_to_csv(path: str, *, max_seconds: float) -> ConvertResult:
    """Real (escaped) CSV output, not the Markdown table pandoc would produce.

    Only the first sheet is written — CSV has no concept of multiple
    sheets, and silently concatenating them would produce invalid CSV.
    """
    if not _has_module("openpyxl"):
        raise UnsupportedFormat(
            "cannot convert .xlsx to CSV — openpyxl is required for CSV "
            "output (pandoc's xlsx support only produces Markdown tables). "
            "Run `pip install gaz[preview]`."
        )
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheet = workbook[workbook.sheetnames[0]]
    buf = io.StringIO()
    writer = csv.writer(buf)
    for row in sheet.iter_rows(values_only=True):
        writer.writerow(["" if v is None else v for v in row])

    warning = None
    if len(workbook.sheetnames) > 1:
        warning = (
            f"workbook has {len(workbook.sheetnames)} sheets; only the "
            f"first ({workbook.sheetnames[0]!r}) was converted to CSV."
        )
    return ConvertResult(text=buf.getvalue(), method="openpyxl", warning=warning)


def _convert_xls(path: str, *, max_seconds: float) -> ConvertResult:
    if shutil.which("pandoc"):
        try:
            proc = _run_subprocess(
                ["pandoc", "--from", "xlsx", "--to", "markdown", path],
                max_seconds=max_seconds,
            )
            if proc.returncode == 0:
                return ConvertResult(
                    text=proc.stdout.decode("utf-8", errors="replace"),
                    method="pandoc",
                )
        except subprocess.TimeoutExpired:
            return ConvertResult(
                text="", method="pandoc", complete=False,
                warning=f"pandoc did not finish within {max_seconds}s",
            )
    raise UnsupportedFormat(
        "cannot convert legacy .xls — pandoc could not read it and no "
        "Python fallback exists for the old binary xls format. Convert it "
        "to .xlsx first, or install pandoc."
    )


def _convert_pdf(path: str, *, max_seconds: float) -> ConvertResult:
    if shutil.which("pdftotext"):
        try:
            proc = _run_subprocess(["pdftotext", path, "-"], max_seconds=max_seconds)
        except subprocess.TimeoutExpired:
            return ConvertResult(
                text="", method="pdftotext", complete=False,
                warning=f"pdftotext did not finish within {max_seconds}s",
            )
        if proc.returncode == 0:
            return ConvertResult(text=proc.stdout.decode("utf-8", errors="replace"), method="pdftotext")

    if _has_module("pypdf"):
        return _convert_pdf_python(path)

    raise UnsupportedFormat(
        "cannot convert .pdf — no converter available. Install poppler "
        "(provides pdftotext), or run `pip install gaz[preview]` for a "
        "Python fallback."
    )


def _convert_pdf_python(path: str) -> ConvertResult:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages_text = [page.extract_text() or "" for page in reader.pages]
    return ConvertResult(text="\n\n".join(pages_text), method="pypdf")


def _pretty_json(path: str) -> ConvertResult:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
    except (OSError, json.JSONDecodeError) as e:
        raise UnsupportedFormat(f"cannot parse {path!r} as JSON: {e}")
    return ConvertResult(text=json.dumps(data, indent=2, sort_keys=False), method="stdlib-json")


def _pretty_yaml(path: str) -> ConvertResult:
    if not _has_module("yaml"):
        warning = "PyYAML not installed — showing raw file content. Run `pip install gaz[preview]`."
        result = _read_as_is(path, method="text (PyYAML not installed)")
        result.warning = warning
        return result
    import yaml

    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = yaml.safe_load(f)
    except (OSError, yaml.YAMLError) as e:
        raise UnsupportedFormat(f"cannot parse {path!r} as YAML: {e}")
    return ConvertResult(text=yaml.dump(data, sort_keys=False, allow_unicode=True), method="pyyaml")


def _pretty_toml(path: str) -> ConvertResult:
    if sys.version_info >= (3, 11):
        import tomllib

        try:
            with open(path, "rb") as f:
                data = tomllib.load(f)
        except (OSError, tomllib.TOMLDecodeError) as e:
            raise UnsupportedFormat(f"cannot parse {path!r} as TOML: {e}")
        return ConvertResult(text=json.dumps(data, indent=2, default=str), method="stdlib-tomllib")

    if not _has_module("tomli"):
        warning = "tomli not installed — showing raw file content. Run `pip install gaz[preview]`."
        result = _read_as_is(path, method="text (tomli not installed)")
        result.warning = warning
        return result
    import tomli

    try:
        with open(path, "rb") as f:
            data = tomli.load(f)
    except (OSError, tomli.TOMLDecodeError) as e:
        raise UnsupportedFormat(f"cannot parse {path!r} as TOML: {e}")
    return ConvertResult(text=json.dumps(data, indent=2, default=str), method="tomli")


def _pretty_xml(path: str) -> ConvertResult:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        dom = xml.dom.minidom.parseString(content)
    except (OSError, xml.parsers.expat.ExpatError) as e:
        raise UnsupportedFormat(f"cannot parse {path!r} as XML: {e}")
    pretty = dom.toprettyxml(indent="  ")
    # minidom emits blank lines for whitespace-only text nodes; drop them.
    lines = [line for line in pretty.splitlines() if line.strip()]
    return ConvertResult(text="\n".join(lines), method="stdlib-xml")


def _pretty_csv(path: str) -> ConvertResult:
    try:
        with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
            rows = list(csv.reader(f))
    except OSError as e:
        raise UnsupportedFormat(f"cannot read {path!r}: {e}")

    if not rows:
        return ConvertResult(text="", method="stdlib-csv")

    widths = [0] * max(len(r) for r in rows)
    for row in rows:
        for i, cell in enumerate(row):
            widths[i] = max(widths[i], len(cell))

    lines = []
    for i, row in enumerate(rows):
        padded = [cell.ljust(widths[j]) for j, cell in enumerate(row)]
        lines.append("  ".join(padded).rstrip())
        if i == 0:
            lines.append("  ".join("-" * w for w in widths))
    return ConvertResult(text="\n".join(lines), method="stdlib-csv")
